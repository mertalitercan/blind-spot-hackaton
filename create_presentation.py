"""Generate BlindSpot hackathon presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Constants
BG = RGBColor(0x05, 0x05, 0x10)
DARK_PANEL = RGBColor(0x0F, 0x0F, 0x1A)
BORDER = RGBColor(0x1E, 0x1E, 0x35)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xF8, 0x71, 0x71)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x64, 0x74, 0x8B)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(text_frame, text, size=14, color=WHITE, bold=False, space_before=0, align=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = Pt(space_before)
    p.alignment = align
    return p


def add_rect(slide, left, top, width, height, fill_color=DARK_PANEL, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

# Accent dots
for i, opacity in enumerate([GREEN, RGBColor(0x2E, 0x7D, 0x32), RGBColor(0x1A, 0x4A, 0x1E)]):
    add_circle(slide, Inches(5.9 + i * 0.45), Inches(2.2), Inches(0.15), opacity)

# Title
add_text(slide, Inches(0), Inches(2.6), SLIDE_W, Inches(1),
         "BlindSpot", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Slogan
add_text(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.6),
         '"We are humans, not machines."', size=20, color=GREEN,
         bold=False, align=PP_ALIGN.CENTER)

# Subtitle
add_text(slide, Inches(0), Inches(4.2), SLIDE_W, Inches(0.5),
         "AI-Powered Fraud Detection System", size=16, color=GRAY,
         align=PP_ALIGN.CENTER)

# Bottom line
add_text(slide, Inches(0), Inches(6.4), SLIDE_W, Inches(0.4),
         "TD Bank  \u00b7  Best AI Hack to Detect Financial Fraud", size=11, color=DARK_GRAY,
         align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 2: The Problem - Fraud Stats
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Problem", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         "Fraud is growing faster than the systems built to detect it.", size=14, color=GRAY)

# Canada stat box
box = add_rect(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(2.6))
txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(3.2), Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CANADA"
p.font.size = Pt(11)
p.font.color.rgb = GREEN
p.font.bold = True
add_para(tf, "$569M", size=44, color=WHITE, bold=True, space_before=6)
add_para(tf, "lost to fraud in 2023", size=13, color=GRAY, space_before=2)
add_para(tf, "+40% increase since 2021", size=11, color=RED, space_before=8)

# Global stat box
box = add_rect(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(2.6))
txBox = slide.shapes.add_textbox(Inches(5.1), Inches(2.0), Inches(3.2), Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "GLOBAL"
p.font.size = Pt(11)
p.font.color.rgb = GREEN
p.font.bold = True
add_para(tf, "$485B", size=44, color=WHITE, bold=True, space_before=6)
add_para(tf, "lost to fraud globally in 2023", size=13, color=GRAY, space_before=2)
add_para(tf, "Projected to exceed $1T by 2027", size=11, color=RED, space_before=8)

# APP fraud box
box = add_rect(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(2.6))
txBox = slide.shapes.add_textbox(Inches(9.1), Inches(2.0), Inches(3.2), Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "APP FRAUD"
p.font.size = Pt(11)
p.font.color.rgb = GREEN
p.font.bold = True
add_para(tf, "70%", size=44, color=WHITE, bold=True, space_before=6)
add_para(tf, "of fraud is Authorized Push Payment", size=13, color=GRAY, space_before=2)
add_para(tf, "Victims willingly send money under coercion", size=11, color=AMBER, space_before=8)

# Bottom insight
add_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2),
         "The biggest blind spot: traditional systems can\u2019t detect fraud when the real account holder is the one pressing \u201cSend.\u201d\n"
         "APP fraud bypasses every rule-based check because the device is real, the IP is real, the credentials are real. The only anomaly is the human.",
         size=13, color=GRAY)


# ═══════════════════════════════════════════════
# SLIDE 3: Traditional vs BlindSpot comparison
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "What Makes Us Different", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         "Our system scrutinizes every layer. Traditional systems only scratch the surface.", size=14, color=GRAY)

# Table headers
cols = [
    ("Capability", Inches(3.5)),
    ("Traditional", Inches(2.2)),
    ("Advanced", Inches(2.2)),
    ("BlindSpot", Inches(2.2)),
]
col_x = [Inches(0.8)]
for i in range(1, len(cols)):
    col_x.append(col_x[-1] + cols[i-1][1])

# Header row
for i, (label, w) in enumerate(cols):
    c = GREEN if i == 3 else GRAY
    bg_c = RGBColor(0x15, 0x25, 0x18) if i == 3 else DARK_PANEL
    r = add_rect(slide, col_x[i], Inches(1.7), w, Inches(0.5), fill_color=bg_c)
    add_text(slide, col_x[i] + Inches(0.15), Inches(1.75), w - Inches(0.3), Inches(0.4),
             label, size=11, color=c, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("Transaction amount & velocity checks", "\u2713", "\u2713", "\u2713"),
    ("Device fingerprinting", "\u2717", "\u2713", "\u2713"),
    ("IP geolocation & impossible travel", "\u2717", "\u2713", "\u2713"),
    ("Recipient network graph analysis", "\u2717", "\u2717", "\u2713"),
    ("Real-time behavioral biometrics", "\u2717", "\u2717", "\u2713"),
    ("Cognitive state detection (coercion/stress)", "\u2717", "\u2717", "\u2713"),
    ("Typing rhythm & hesitation analysis", "\u2717", "\u2717", "\u2713"),
    ("AI-driven reasoning (not just rules)", "\u2717", "\u2717", "\u2713"),
    ("Dynamic scoring that improves over time", "\u2717", "\u2717", "\u2713"),
]

for row_i, (cap, trad, adv, bs) in enumerate(rows):
    y = Inches(2.25) + Inches(row_i * 0.48)
    # Capability
    add_rect(slide, col_x[0], y, cols[0][1], Inches(0.44), fill_color=DARK_PANEL if row_i % 2 == 0 else BG)
    add_text(slide, col_x[0] + Inches(0.15), y + Inches(0.05), cols[0][1] - Inches(0.3), Inches(0.35),
             cap, size=10, color=WHITE)
    # Values
    for ci, val in enumerate([trad, adv, bs], 1):
        bg_c = DARK_PANEL if row_i % 2 == 0 else BG
        if ci == 3:
            bg_c = RGBColor(0x12, 0x20, 0x15) if row_i % 2 == 0 else RGBColor(0x0A, 0x15, 0x0D)
        add_rect(slide, col_x[ci], y, cols[ci][1], Inches(0.44), fill_color=bg_c)
        vc = GREEN if val == "\u2713" else RGBColor(0x40, 0x40, 0x50)
        add_text(slide, col_x[ci], y + Inches(0.05), cols[ci][1], Inches(0.35),
                 val, size=14, color=vc, bold=True, align=PP_ALIGN.CENTER)

# Bottom callout
add_rect(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), fill_color=RGBColor(0x15, 0x25, 0x18), border_color=GREEN)
add_text(slide, Inches(1.0), Inches(6.75), Inches(11.3), Inches(0.4),
         "BlindSpot doesn\u2019t just flag numbers \u2014 it deeply scrutinizes the human behind every transaction.", size=12, color=GREEN,
         bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 4: Data We Collect (with code)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "What We Capture", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         "Every interaction in the mobile app is a data point that builds a behavioral fingerprint.", size=14, color=GRAY)

# Left: categories
categories = [
    ("Typing Biometrics", [
        "typing_speed_wpm",
        "typing_rhythm_signature",
        "error_rate",
        "segmented_typing_detected",
    ]),
    ("Touch Biometrics", [
        "avg_touch_pressure",
        "avg_touch_radius",
        "hand_dominance",
    ]),
    ("Behavioral Signals", [
        "paste_detected  /  paste_field",
        "confirm_button_hesitation_ms",
        "confirm_attempts",
        "navigation_directness_score",
        "screen_familiarity_score",
    ]),
    ("Session Context", [
        "is_phone_call_active",
        "total_dead_time_ms",
        "app_switches",
    ]),
]

y_pos = Inches(1.7)
for cat_name, fields in categories:
    add_text(slide, Inches(0.8), y_pos, Inches(4), Inches(0.3),
             cat_name, size=12, color=GREEN, bold=True)
    y_pos += Inches(0.3)
    for field in fields:
        add_text(slide, Inches(1.0), y_pos, Inches(4), Inches(0.25),
                 field, size=10, color=GRAY, font_name="Consolas")
        y_pos += Inches(0.22)
    y_pos += Inches(0.12)

# Right: code block
code_bg = add_rect(slide, Inches(5.5), Inches(1.7), Inches(7.0), Inches(5.3),
                   fill_color=RGBColor(0x0A, 0x0A, 0x12), border_color=BORDER)

code_text = """export interface TelemetrySnapshot {
  session_id: string;
  user_id: string;
  session_duration_ms: number;

  // Typing biometrics
  keystroke_events: KeystrokeEvent[];
  typing_speed_wpm: number;
  error_rate: number;
  typing_rhythm_signature: number[];
  segmented_typing_detected: boolean;

  // Paste & hesitation
  paste_detected: boolean;
  paste_field: string;
  confirm_button_hesitation_ms: number;
  confirm_attempts: number;

  // Touch & navigation
  touch_events: TouchEvent[];
  avg_touch_pressure: number;
  hand_dominance: string;
  navigation_directness_score: number;
  screen_familiarity_score: number;

  // Dead time & app switches
  total_dead_time_ms: number;
  dead_time_periods: DeadTimePeriod[];
  app_switches: AppSwitch[];
}"""

txBox = slide.shapes.add_textbox(Inches(5.7), Inches(1.85), Inches(6.6), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0xA0, 0xC0, 0xA0)
p.font.name = "Consolas"
p.line_spacing = Pt(13)

# file path label
add_text(slide, Inches(5.5), Inches(7.05), Inches(7.0), Inches(0.3),
         "mobile/lib/telemetry.ts", size=9, color=DARK_GRAY, align=PP_ALIGN.RIGHT,
         font_name="Consolas")


# ═══════════════════════════════════════════════
# SLIDE 5: Multi-Agent System
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Multi-Agent AI System", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         "6 specialized AI agents analyze every transaction in parallel, then a meta-scorer synthesizes the verdict.", size=14, color=GRAY)

agents = [
    ("Behavioral", "Compares typing speed, rhythm, pressure, and navigation patterns against the user's historical baseline.", "20%"),
    ("Cognitive", "Detects coercion, stress, and coaching signals — phone calls, hesitation, dead time, segmented typing.", "30%"),
    ("Transaction", "Analyzes amount, recipient, MCC codes, geo-location, and patterns vs. user's income and history.", "20%"),
    ("Device", "Verifies device fingerprint, IP geolocation, VPN/emulator detection, impossible travel scenarios.", "10%"),
    ("Graph", "Maps recipient network — fan-in patterns, account age, first-time senders, mule indicators.", "20%"),
]

# Agent cards in 2 rows
for i, (name, desc, weight) in enumerate(agents):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.0)
    y = Inches(1.7) + Inches(row * 2.4)

    card = add_rect(slide, x, y, Inches(3.7), Inches(2.1))

    # Green dot + name
    add_circle(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.1), GREEN)
    add_text(slide, x + Inches(0.4), y + Inches(0.15), Inches(2.5), Inches(0.3),
             name, size=15, color=WHITE, bold=True)

    # Weight badge
    add_text(slide, x + Inches(2.7), y + Inches(0.15), Inches(0.8), Inches(0.3),
             weight, size=12, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)

    # Description
    add_text(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.3), Inches(1.3),
             desc, size=11, color=GRAY)

# Meta-scorer at bottom
meta_y = Inches(1.7) + Inches(2.4)
meta_x = Inches(4.8)
card = add_rect(slide, meta_x, meta_y, Inches(7.7), Inches(2.1),
                fill_color=RGBColor(0x15, 0x25, 0x18), border_color=GREEN)
add_circle(slide, meta_x + Inches(0.2), meta_y + Inches(0.25), Inches(0.1), GREEN)
add_text(slide, meta_x + Inches(0.4), meta_y + Inches(0.15), Inches(4), Inches(0.3),
         "Meta-Reasoning Scorer", size=15, color=GREEN, bold=True)
add_text(slide, meta_x + Inches(0.2), meta_y + Inches(0.55), Inches(7.3), Inches(1.4),
         "Synthesizes all 5 agent outputs using weighted scoring + amplification rules.\n"
         "Applies AI reasoning to produce final fraud score, risk level, and recommended actions.\n"
         "Outputs a human-readable explanation of why the transaction was flagged.",
         size=11, color=GRAY)

# Arrow label
add_text(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         "All 6 agents run concurrently  \u2192  Results streamed in real-time to the analyst dashboard",
         size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 6: Scoring System
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Dynamic Scoring System", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Weighted formula + AI amplification rules. The system learns and adapts with every transaction.", size=14, color=GRAY)

# Formula box
formula_box = add_rect(slide, Inches(0.8), Inches(1.7), Inches(7.5), Inches(1.6),
                        fill_color=RGBColor(0x0A, 0x0A, 0x12))
formula = (
    "score = Behavioral \u00d7 0.20  +  Cognitive \u00d7 0.30  +  Transaction \u00d7 0.20\n"
    "            + Device \u00d7 0.10  +  Graph \u00d7 0.20"
)
add_text(slide, Inches(1.0), Inches(1.9), Inches(7.1), Inches(0.8),
         formula, size=14, color=WHITE, bold=True, font_name="Consolas")

add_text(slide, Inches(1.0), Inches(2.7), Inches(7.1), Inches(0.4),
         "Cognitive has the highest weight (30%) \u2014 because detecting human coercion is our core differentiator.", size=11, color=GREEN)

# Amplification rules
add_text(slide, Inches(0.8), Inches(3.6), Inches(7.5), Inches(0.4),
         "Amplification Rules", size=16, color=WHITE, bold=True)

rules = [
    ("APP Fraud Pattern", "Cognitive > 70 AND Device < 50", "\u00d71.4", "Real device, coerced user"),
    ("Stressed Behavior", "Cognitive > 70 AND Behavioral > 70", "\u00d71.3", "Unusual behavior under stress"),
    ("Mule Network", "Graph > 70 AND Transaction > 60", "\u00d71.3", "Suspicious recipient + unusual txn"),
    ("Account Takeover", "Behavioral > 80 AND Device > 80", "\u00d71.4", "Different person, different device"),
]

for i, (name, condition, mult, reason) in enumerate(rules):
    y = Inches(4.1) + Inches(i * 0.55)
    add_rect(slide, Inches(0.8), y, Inches(7.5), Inches(0.48), fill_color=DARK_PANEL if i % 2 == 0 else BG)
    add_text(slide, Inches(1.0), y + Inches(0.06), Inches(1.8), Inches(0.35),
             name, size=10, color=WHITE, bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.06), Inches(3.0), Inches(0.35),
             condition, size=9, color=GRAY, font_name="Consolas")
    add_text(slide, Inches(5.8), y + Inches(0.06), Inches(0.6), Inches(0.35),
             mult, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.4), y + Inches(0.06), Inches(1.8), Inches(0.35),
             reason, size=9, color=DARK_GRAY)

# Risk levels on the right
add_text(slide, Inches(8.8), Inches(1.7), Inches(3.7), Inches(0.4),
         "Risk Levels", size=16, color=WHITE, bold=True)

levels = [
    ("0\u201330", "Low Risk", GREEN, "Allow transaction"),
    ("31\u201370", "Medium Risk", AMBER, "Flag for review"),
    ("71\u2013100", "High Risk", RED, "Block & pause account"),
]

for i, (range_s, label, color, action) in enumerate(levels):
    y = Inches(2.3) + Inches(i * 0.9)
    add_rect(slide, Inches(8.8), y, Inches(3.7), Inches(0.75), fill_color=DARK_PANEL)
    add_text(slide, Inches(9.0), y + Inches(0.08), Inches(1.0), Inches(0.3),
             range_s, size=20, color=color, bold=True, font_name="Consolas")
    add_text(slide, Inches(10.0), y + Inches(0.08), Inches(2.3), Inches(0.3),
             label, size=13, color=color, bold=True)
    add_text(slide, Inches(10.0), y + Inches(0.38), Inches(2.3), Inches(0.3),
             action, size=10, color=GRAY)

# Dynamic callout
add_rect(slide, Inches(8.8), Inches(5.1), Inches(3.7), Inches(1.2),
         fill_color=RGBColor(0x15, 0x25, 0x18), border_color=GREEN)
add_text(slide, Inches(9.0), Inches(5.2), Inches(3.3), Inches(1.0),
         "Self-Improving System\n\n"
         "Every transaction and behavioral input expands the user's baseline. "
         "The more the system sees, the better it gets at distinguishing real users from imposters.",
         size=10, color=GREEN)


# ═══════════════════════════════════════════════
# SLIDE 7: Demo Video
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(0.8),
         "Live Demo", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Play button circle
play = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(3.6), Inches(1.5), Inches(1.5))
play.fill.solid()
play.fill.fore_color.rgb = GREEN
play.line.fill.background()

add_text(slide, Inches(5.9), Inches(3.85), Inches(1.5), Inches(1.0),
         "\u25B6", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.4), SLIDE_W, Inches(0.4),
         "Mobile App  \u2192  Real-time Agent Pipeline  \u2192  Analyst Dashboard  \u2192  AI Chat", size=13, color=GRAY,
         align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.9), SLIDE_W, Inches(0.3),
         "Insert 1-minute demo video here", size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 8: Tech Stack
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0), Inches(0.5), SLIDE_W, Inches(0.6),
         "Tech Stack", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

stack_sections = [
    ("Backend", [("FastAPI", "Python"), ("Uvicorn", "ASGI"), ("Pydantic", "Validation")]),
    ("AI & Agents", [("Claude AI", "Anthropic"), ("Railtracks", "Agent Framework"), ("Claude Haiku 4.5", "Agents"), ("Claude Sonnet 4", "Chat")]),
    ("Frontend", [("Next.js 15", "React"), ("Tailwind CSS", "Styling"), ("WebSocket", "Real-time")]),
    ("Mobile", [("Expo", "Framework"), ("React Native", "UI"), ("Custom Telemetry", "Biometrics")]),
    ("Database", [("Supabase", "PostgreSQL")]),
]

section_x = Inches(0.5)
for sec_i, (section_name, items) in enumerate(stack_sections):
    x = section_x + Inches(sec_i * 2.5)
    y_start = Inches(1.5)

    add_text(slide, x, y_start, Inches(2.3), Inches(0.4),
             section_name, size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    for item_i, (name, sub) in enumerate(items):
        y = y_start + Inches(0.6) + Inches(item_i * 1.1)
        card = add_rect(slide, x, y, Inches(2.3), Inches(0.9))
        add_text(slide, x, y + Inches(0.12), Inches(2.3), Inches(0.35),
                 name, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(0.48), Inches(2.3), Inches(0.3),
                 sub, size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 9: Thank You
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

for i, opacity in enumerate([GREEN, RGBColor(0x2E, 0x7D, 0x32), RGBColor(0x1A, 0x4A, 0x1E)]):
    add_circle(slide, Inches(5.9 + i * 0.45), Inches(2.4), Inches(0.15), opacity)

add_text(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(0.8),
         "Thank You", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3.7), SLIDE_W, Inches(0.5),
         '"We are humans, not machines."', size=18, color=GREEN, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(4.8), SLIDE_W, Inches(0.4),
         "BlindSpot  \u00b7  AI-Powered Fraud Detection", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.4),
         "Questions?", size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
output_path = "/Users/mali/Desktop/projects/blind-spot/BlindSpot_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
